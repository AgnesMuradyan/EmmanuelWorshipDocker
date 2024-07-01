from django.db import models
from pptx import Presentation
import io
import time
import io
from django.db import models
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import logging

logger = logging.getLogger(__name__)

class Album(models.Model):
    id = models.AutoField(primary_key=True)
    title = models.CharField(max_length=200)

    def __str__(self):
        return self.title

class Song(models.Model):
    id = models.AutoField(primary_key=True)
    title = models.CharField(max_length=200)
    original_link = models.URLField(max_length=200, blank=True, null=True)
    original_key = models.CharField(max_length=200, blank=True, null=True)
    album = models.ForeignKey(Album, related_name='songs', on_delete=models.CASCADE)
    verse = models.TextField(blank=True, null=True)
    chords = models.BinaryField(null=True, blank=True)
    powerpoint = models.BinaryField(null=True, blank=True)
    created_at = models.DateTimeField(auto_now_add=True)
    updated_at = models.DateTimeField(auto_now=True)

    def __str__(self):
        return self.title


class Member(models.Model):
    id = models.AutoField(primary_key=True)
    first_name = models.CharField(max_length=100)
    last_name = models.CharField(max_length=100)
    birth_date = models.DateField(null=True, blank=True)
    phone_number = models.CharField(max_length=20, null=True, blank=True)

    class Meta:
        abstract = True

    def __str__(self):
        return f"{self.first_name} {self.last_name}"


class Instrument(models.Model):
    name = models.CharField(max_length=100)

    def __str__(self):
        return self.name


class Musician(Member):
    instruments = models.ManyToManyField(Instrument, through='MusicianInstrument', related_name='musicians')


class MusicianInstrument(models.Model):
    musician = models.ForeignKey(Musician, related_name='musician_instruments', on_delete=models.CASCADE)
    instrument = models.ForeignKey(Instrument, related_name='musician_instruments', on_delete=models.CASCADE)

    def __str__(self):
        return f"{self.musician} plays {self.instrument}"


class Singer(Member):
    SOLOIST = 'SO'
    CHOIR = 'CH'
    ROLE_CHOICES = [
        (SOLOIST, 'Soloist'),
        (CHOIR, 'Choir'),
    ]
    role = models.CharField(max_length=2, choices=ROLE_CHOICES)

    def __str__(self):
        role_str = "Soloist" if self.role == self.SOLOIST else "Choir"
        return f"{self.first_name} {self.last_name} - {role_str}"


# class Plan(models.Model):
#     date = models.DateField()
#     concatenated_powerpoint = models.BinaryField(null=True, blank=True)
#
#     lead_singers = models.ManyToManyField('Singer', related_name='lead_plans')
#     singers = models.ManyToManyField('Singer', related_name='plans')
#     musicians = models.ManyToManyField('Musician', related_name='plans')
#     songs = models.ManyToManyField('Song', related_name='plans')
#
#     def __str__(self):
#         return f"Plan for {self.date}"
#
#     def save(self, *args, **kwargs):
#         # Save the Plan instance to generate the ID if new
#         super().save(*args, **kwargs)
#         # Ensure songs are added before creating the PowerPoint
#         if not self.concatenated_powerpoint:
#             self.create_concatenated_powerpoint()
#             super().save(update_fields=['concatenated_powerpoint'])  # Save only the concatenated_powerpoint field
#
#     def create_concatenated_powerpoint(self):
#         print("Creating new PowerPoint presentation")
#         prs = Presentation()
#         for song in self.songs.all():
#             print(f"Processing song: {song.title}")
#             if song.powerpoint:
#                 print(f"Adding slides from song: {song.title}")
#                 song_prs = Presentation(io.BytesIO(song.powerpoint))
#                 for slide in song_prs.slides:
#                     slide_layout = prs.slide_layouts[5]  # Choose a slide layout
#                     new_slide = prs.slides.add_slide(slide_layout)
#                     for shape in slide.shapes:
#                         if shape.has_text_frame:
#                             new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
#                             new_shape.text = shape.text
#         # Save the concatenated presentation to a byte stream
#         output = io.BytesIO()
#         prs.save(output)
#         output.seek(0)
#         print("PowerPoint presentation created successfully")
#         self.concatenated_powerpoint = output.read()

class Plan(models.Model):
    THURSDAY = 'TH'
    SUNDAY = 'SU'
    OTHER = 'OT'

    DAY_TYPE_CHOICES = [
        (THURSDAY, 'Thursday'),
        (SUNDAY, 'Sunday'),
        (OTHER, 'Other'),
    ]
    date = models.DateField()
    concatenated_powerpoint = models.BinaryField(null=True, blank=True)
    day_type = models.CharField(max_length=2, choices=DAY_TYPE_CHOICES, default=OTHER)

    lead_singers = models.ManyToManyField('Singer', related_name='lead_plans', null=True, blank=True)
    singers = models.ManyToManyField('Singer', related_name='plans', null=True, blank=True)
    musicians = models.ManyToManyField('Musician', related_name='plans', null=True, blank=True)
    songs = models.ManyToManyField(Song, through='PlanSong', related_name='plans', null=True, blank=True)

    def __str__(self):
        return f"Plan for {self.date}"

    def create_concatenated_powerpoint(self):
        logger.debug("Creating new PowerPoint presentation")
        logger.info("Creating new PowerPoint presentation")
        print("Creating new PowerPoint presentation")
        prs = Presentation()

        prs.slide_width = Inches(23.00)  # Approx 51.8 cm
        prs.slide_height = Inches(12.00)  # Approx 30.0 cm

        plan_songs = self.plansong_set.all().order_by('order')

        if not plan_songs:
            print("No PlanSongs found.")
            return None

        for plan_song in plan_songs:
            song = plan_song.song
            print(f"Processing song: {song.title}")
            if song.powerpoint:
                try:
                    song_prs = Presentation(io.BytesIO(song.powerpoint))
                except Exception as e:
                    print(f"Error reading PowerPoint for song {song.title}: {e}")
                    continue

                print(f"Adding slides from song: {song.title}")
                for slide in song_prs.slides:
                    slide_layout = prs.slide_layouts[5]  # Choose a slide layout
                    new_slide = prs.slides.add_slide(slide_layout)

                    # Set slide background to black
                    background = new_slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(0, 0, 0)

                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                            text_frame = new_shape.text_frame
                            text_frame.text = shape.text

                            # Center the text and set font properties
                            for paragraph in text_frame.paragraphs:
                                paragraph.alignment = PP_ALIGN.CENTER
                                for run in paragraph.runs:
                                    run.font.size = Pt(81)
                                    run.font.name = 'Agg_Book1'
                                    run.font.bold = True
                                    run.font.color.rgb = RGBColor(255, 255, 255)

        # Save the concatenated presentation to a byte stream
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        print("PowerPoint presentation created successfully")
        return output.read()


class PlanSong(models.Model):
    plan = models.ForeignKey(Plan, on_delete=models.CASCADE)
    song = models.ForeignKey(Song, on_delete=models.CASCADE)
    order = models.PositiveIntegerField()

    class Meta:
        unique_together = ('plan', 'order')  # Ensure each order is unique per plan
        ordering = ['order']

    # def __str__(self):
    #     return f"{self.plan} - {self.song} (Order: {self.order})"
